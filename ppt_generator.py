import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from datetime import datetime

class PPTGenerator:
    def __init__(self):
        self.prs = None
    
    def create_presentation(self, title, subtitle=""):
        self.prs = Presentation()
        self.prs.slide_width = Inches(16)
        self.prs.slide_height = Inches(9)
        
        slide_layout = self.prs.slide_layouts[0]
        slide = self.prs.slides.add_slide(slide_layout)
        
        title_shape = slide.shapes.title
        title_shape.text = title
        title_shape.text_frame.paragraphs[0].font.size = Pt(36)
        title_shape.text_frame.paragraphs[0].font.bold = True
        title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
        
        if subtitle:
            subtitle_shape = slide.placeholders[1]
            subtitle_shape.text = subtitle
            subtitle_shape.text_frame.paragraphs[0].font.size = Pt(18)
        
        return self.prs
    
    def add_title_slide(self, title, subtitle=""):
        slide_layout = self.prs.slide_layouts[0]
        slide = self.prs.slides.add_slide(slide_layout)
        
        title_shape = slide.shapes.title
        title_shape.text = title
        title_shape.text_frame.paragraphs[0].font.size = Pt(32)
        title_shape.text_frame.paragraphs[0].font.bold = True
        title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
        
        if subtitle:
            subtitle_shape = slide.placeholders[1]
            subtitle_shape.text = subtitle
            subtitle_shape.text_frame.paragraphs[0].font.size = Pt(16)
        
        return slide
    
    def add_content_slide(self, title, content_items, layout_type='bullets'):
        if layout_type == 'bullets':
            slide_layout = self.prs.slide_layouts[1]
            slide = self.prs.slides.add_slide(slide_layout)
            
            title_shape = slide.shapes.title
            title_shape.text = title
            
            content_shape = slide.placeholders[1]
            tf = content_shape.text_frame
            tf.clear()
            
            for item in content_items:
                p = tf.add_paragraph()
                p.text = item
                p.level = 0
                p.font.size = Pt(18)
        
        elif layout_type == 'two_column':
            slide_layout = self.prs.slide_layouts[3]
            slide = self.prs.slides.add_slide(slide_layout)
            
            title_shape = slide.shapes.title
            title_shape.text = title
            
            left_content = content_items.get('left', [])
            right_content = content_items.get('right', [])
            
            left_placeholder = slide.placeholders[1]
            left_tf = left_placeholder.text_frame
            left_tf.clear()
            for item in left_content:
                p = left_tf.add_paragraph()
                p.text = item
                p.font.size = Pt(16)
            
            right_placeholder = slide.placeholders[2]
            right_tf = right_placeholder.text_frame
            right_tf.clear()
            for item in right_content:
                p = right_tf.add_paragraph()
                p.text = item
                p.font.size = Pt(16)
        
        elif layout_type == 'title_only':
            slide_layout = self.prs.slide_layouts[5]
            slide = self.prs.slides.add_slide(slide_layout)
            
            title_shape = slide.shapes.title
            title_shape.text = title
        
        return slide
    
    def add_table_slide(self, title, headers, data):
        slide_layout = self.prs.slide_layouts[5]
        slide = self.prs.slides.add_slide(slide_layout)
        
        title_shape = slide.shapes.title
        title_shape.text = title
        
        rows = len(data) + 1
        cols = len(headers)
        
        left = Inches(0.5)
        top = Inches(1.5)
        width = Inches(15)
        height = Inches(0.8 * rows)
        
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table
        
        for i, header in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = header
            cell.text_frame.paragraphs[0].font.bold = True
            cell.text_frame.paragraphs[0].font.size = Pt(14)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0, 51, 102)
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        
        for row_idx, row_data in enumerate(data):
            for col_idx, value in enumerate(row_data):
                cell = table.cell(row_idx + 1, col_idx)
                cell.text = str(value)
                cell.text_frame.paragraphs[0].font.size = Pt(12)
        
        return slide
    
    def add_chart_slide(self, title, chart_data):
        slide_layout = self.prs.slide_layouts[5]
        slide = self.prs.slides.add_slide(slide_layout)
        
        title_shape = slide.shapes.title
        title_shape.text = title
        
        left = Inches(1)
        top = Inches(1.5)
        width = Inches(14)
        height = Inches(6)
        
        categories = chart_data.get('categories', [])
        values = chart_data.get('values', [])
        series_name = chart_data.get('series_name', '数据')
        
        chart = slide.shapes.add_chart(
            chart_data.get('chart_type', 'column'),
            left, top, width, height
        ).chart
        
        chart.series[0].name = series_name
        chart.series[0].values = values
        chart.category_labels = categories
        
        chart.has_title = True
        chart.chart_title.text_frame.text = chart_data.get('chart_title', '')
        
        return slide
    
    def add_summary_slide(self, summary_items):
        slide_layout = self.prs.slide_layouts[5]
        slide = self.prs.slides.add_slide(slide_layout)
        
        title_shape = slide.shapes.title
        title_shape.text = "总结"
        
        content = ""
        for key, value in summary_items.items():
            content += f"● {key}: {value}\n"
        
        left = Inches(1)
        top = Inches(2)
        width = Inches(14)
        height = Inches(5)
        
        textbox = slide.shapes.add_textbox(left, top, width, height)
        tf = textbox.text_frame
        tf.text = content
        
        for paragraph in tf.paragraphs:
            paragraph.font.size = Pt(18)
        
        return slide
    
    def add_meeting_minutes_slides(self, minutes_data):
        self.add_title_slide("会议纪要", minutes_data.get('date', ''))
        
        self.add_table_slide("会议信息", 
                           ['项目', '内容'],
                           [
                               ['会议时间', minutes_data.get('date', '') + ' ' + minutes_data.get('time', '')],
                               ['会议地点', minutes_data.get('location', '未指定')],
                               ['参会人数', str(len(minutes_data.get('participants', []))) + '人'],
                               ['缺席人数', str(len(minutes_data.get('absentees', []))) + '人']
                           ])
        
        if minutes_data.get('agenda'):
            self.add_content_slide("会议议程", minutes_data['agenda'])
        
        if minutes_data.get('decisions'):
            self.add_content_slide("会议决议", minutes_data['decisions'])
        
        if minutes_data.get('action_items'):
            self.add_content_slide("行动项", minutes_data['action_items'])
        
        if minutes_data.get('summary'):
            self.add_content_slide("会议总结", [minutes_data['summary']])
    
    def add_attendance_slides(self, attendance_data):
        self.add_title_slide("签到统计报告", attendance_data.get('date', ''))
        
        rate = (attendance_data.get('total_attended', 0) / attendance_data.get('total_expected', 1) * 100)
        
        self.add_table_slide("统计信息",
                           ['项目', '数据'],
                           [
                               ['统计日期', attendance_data.get('date', '')],
                               ['应到人数', str(attendance_data.get('total_expected', 0))],
                               ['实到人数', str(attendance_data.get('total_attended', 0))],
                               ['签到率', f"{rate:.1f}%"]
                           ])
        
        if attendance_data.get('attendees'):
            attendees_list = [f"{i+1}. {a['name']}" for i, a in enumerate(attendance_data['attendees'][:10])]
            if len(attendance_data['attendees']) > 10:
                attendees_list.append(f"... 等{len(attendance_data['attendees'])}人")
            self.add_content_slide("签到人员", attendees_list)
        
        if attendance_data.get('absentees'):
            absentees_list = [f"{i+1}. {a['name']}" for i, a in enumerate(attendance_data['absentees'])]
            self.add_content_slide("缺席人员", absentees_list)
    
    def add_training_slides(self, training_data):
        self.add_title_slide("培训效果报告", training_data.get('date', ''))
        
        self.add_table_slide("培训信息",
                           ['项目', '内容'],
                           [
                               ['培训名称', training_data.get('title', '未指定')],
                               ['培训日期', training_data.get('date', '')],
                               ['培训师', training_data.get('trainer', '未指定')],
                               ['参与人数', str(training_data.get('total_participants', 0))],
                               ['平均评分', f"{training_data.get('overall_score', 0):.1f}分"]
                           ])
        
        if training_data.get('comments'):
            comments_list = [f"{c['name']}: {c['comment']}" for c in training_data['comments'][:5]]
            if len(training_data['comments']) > 5:
                comments_list.append(f"... 共{len(training_data['comments'])}条反馈")
            self.add_content_slide("学员反馈", comments_list)
    
    def generate_monthly_report(self, data, output_path):
        month = datetime.now().strftime('%Y年%m月')
        self.create_presentation(f"{month}月度汇报材料", "珠海华润银行大厦")
        
        if data.get('meeting_minutes'):
            self.add_meeting_minutes_slides(data['meeting_minutes'])
        
        if data.get('attendance'):
            self.add_attendance_slides(data['attendance'])
        
        if data.get('training'):
            self.add_training_slides(data['training'])
        
        summary_items = {}
        if data.get('meeting_minutes'):
            summary_items['会议次数'] = str(len(data['meeting_minutes'].get('decisions', [])) + 1) + '次'
        if data.get('attendance'):
            rate = (data['attendance'].get('total_attended', 0) / data['attendance'].get('total_expected', 1) * 100)
            summary_items['平均签到率'] = f"{rate:.1f}%"
        if data.get('training'):
            summary_items['培训场次'] = '1场'
            summary_items['培训平均评分'] = f"{data['training'].get('overall_score', 0):.1f}分"
        
        if summary_items:
            self.add_summary_slide(summary_items)
        
        self.prs.save(output_path)
        return output_path

if __name__ == "__main__":
    generator = PPTGenerator()
    
    sample_data = {
        'meeting_minutes': {
            'title': '4月份EHS会议纪要',
            'date': '2026年4月15日',
            'time': '14:00-15:30',
            'location': '会议室A',
            'participants': ['张三', '李四', '王五', '赵六'],
            'absentees': ['钱七'],
            'agenda': ['安全检查情况汇报', '培训计划讨论', '应急预案修订'],
            'decisions': ['每月进行一次全面安全检查', '下月组织消防培训', '更新应急预案'],
            'action_items': ['张三负责更新检查台账', '李四安排培训场地', '王五修订应急预案'],
            'summary': '本次会议总结了本月安全工作情况，明确了下月工作计划。'
        },
        'attendance': {
            'date': '2026年4月',
            'total_expected': 30,
            'total_attended': 28,
            'attendees': [{'name': '张三'}, {'name': '李四'}],
            'absentees': [{'name': '钱七'}]
        },
        'training': {
            'title': '消防安全培训',
            'date': '2026年4月20日',
            'trainer': '李教官',
            'total_participants': 25,
            'overall_score': 92.5,
            'comments': [
                {'name': '张三', 'comment': '培训内容实用'},
                {'name': '李四', 'comment': '消防器材使用讲解清晰'}
            ]
        }
    }
    
    output_path = os.path.join(os.path.dirname(__file__), 'monthly_report.pptx')
    generator.generate_monthly_report(sample_data, output_path)
    print(f"PPT已生成: {output_path}")